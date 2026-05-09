#!/usr/bin/env python3
"""
Generates docs/decks/nee-forge-pitch.pptx from the v2 pitch-deck-plan.md.

Run from repo root:
    python3 docs/decks/build-pitch-deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ─── Brand tokens (match the app design system) ──────────────────────────
INK        = RGBColor(0x0A, 0x0E, 0x0F)
INK_2      = RGBColor(0x1A, 0x1F, 0x21)
INK_3      = RGBColor(0x6B, 0x88, 0xA4)
PAPER      = RGBColor(0xF5, 0xF2, 0xEA)
PAPER_2    = RGBColor(0xEB, 0xE7, 0xDC)
PAPER_3    = RGBColor(0xDD, 0xD7, 0xC8)
ACCENT     = RGBColor(0xC2, 0x41, 0x0C)
WARN       = RGBColor(0xB4, 0x53, 0x09)

DISPLAY_FONT = "Georgia"
SANS_FONT    = "Helvetica Neue"
MONO_FONT    = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT     = Path(__file__).resolve().parent.parent.parent
ASSETS   = ROOT / "docs" / "decks" / "assets"
OUT_FILE = ROOT / "docs" / "decks" / "nee-forge-pitch.pptx"

def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, color)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *,
             font=SANS_FONT, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb

def add_caption(slide, x, y, w, text, size=14):
    return add_text(slide, x, y, w, Inches(0.6), text,
                    font=SANS_FONT, size=size, color=INK_3)

def add_bullets(slide, x, y, w, h, lines, *, font=SANS_FONT, size=18,
                color=INK, line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•   " + line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def add_table(slide, x, y, w, h, headers, rows, *,
              header_color=INK, alt_row=PAPER_2, font=SANS_FONT, size=12):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    tbl_shape = slide.shapes.add_table(rows_n, cols_n, x, y, w, h)
    tbl = tbl_shape.table
    for c, txt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = txt
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = PAPER
    for ri, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if ri % 2 else alt_row
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.word_wrap = True
            for li, line in enumerate(str(txt).split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = line
                r.font.name = font
                r.font.size = Pt(size)
                r.font.color.rgb = INK_2
    return tbl

def slide_bg(slide, color=PAPER):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    bg.shadow.inherit = False
    return bg

def add_footer(slide, num_label):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.35),
             "NEE FORGE  ·  Production Software for the Workshop",
             font=MONO_FONT, size=9, color=INK_3)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
             num_label, font=MONO_FONT, size=9, color=INK_3,
             align=PP_ALIGN.RIGHT)

def header_band(slide, title, kicker=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.4), kicker,
                 font=MONO_FONT, size=11, color=ACCENT, bold=True)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12.4), Inches(0.85), title,
             font=DISPLAY_FONT, size=32, color=INK)


prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── Slide 1 — Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s, INK)
add_text(s, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
         "NEE", font=MONO_FONT, size=18, color=PAPER_3)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(2.2),
         "Forge", font=DISPLAY_FONT, size=130, color=PAPER)
add_text(s, Inches(0.8), Inches(3.7), Inches(11), Inches(0.6),
         "Production Software for the Workshop",
         font=SANS_FONT, size=24, color=PAPER_2)
add_rect(s, Inches(0.8), Inches(5.9), Inches(0.06), Inches(0.5), ACCENT)
add_text(s, Inches(1.05), Inches(5.95), Inches(11), Inches(0.5),
         "Where shop-floor work gets shaped.",
         font=DISPLAY_FONT, size=18, color=PAPER_2)


# ── Slide 2 — The pain today ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "The Monday-morning printout is stale by 11 a.m.", kicker="THE PAIN TODAY")
bullets = [
    '"Where\'s the chassis for Job Card 0007?" → 3 phone calls + a walk to the floor.',
    "Variance is reconciled in arrears. The supervisor learns next Friday why something ran late, not in time to act.",
    "Three or more versions of \"the master sheet\" sit in different inboxes. None is current.",
    "When a tech reports a blocker, it lives in his head until end of day. Sometimes longer.",
    "Every printout is a snapshot of a moment that has already passed.",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5),
            bullets, size=17, line_spacing=1.5)
add_caption(s, Inches(0.8), Inches(6.5), Inches(11.7),
            "Show of hands — who has been on a phone call this week to find out where a job is?",
            size=13)
add_footer(s, "02 / 11")


# ── Slide 3 — NEE Forge in one diagram ───────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "One system. Five personas. One source of truth.",
            kicker="NEE FORGE")
labels = [
    ("SALES", "PDF intake auto-builds Job Card"),
    ("DRAFTER", "Drawings · BOM · Layout upload"),
    ("SUPERVISOR", "Schedule + 4-week capacity heatmap"),
    ("TECH", "Clock in / out from a phone"),
    ("QC", "6-item checklist + customer email"),
]
n = len(labels)
margin_x = Inches(0.4)
gap      = Inches(0.15)
total_w  = SLIDE_W - margin_x * 2
arrow_w  = Inches(0.35)
box_w    = Emu((total_w - arrow_w * (n - 1) - gap * 2 * (n - 1)) // n)
box_h    = Inches(2.2)
top      = Inches(2.5)
x = margin_x
for i, (head, body) in enumerate(labels):
    add_rect(s, x, top, box_w, box_h, PAPER_2)
    add_rect(s, x, top, box_w, Inches(0.06), ACCENT)
    add_text(s, x + Inches(0.15), top + Inches(0.25),
             Emu(box_w - Inches(0.3)), Inches(0.5), head,
             font=MONO_FONT, size=11, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.15), top + Inches(0.85),
             Emu(box_w - Inches(0.3)), Inches(1.2), body,
             font=DISPLAY_FONT, size=17, color=INK)
    x = Emu(x + box_w + gap)
    if i < n - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, top + Inches(0.85),
                                   arrow_w, Inches(0.5))
        fill(arrow, INK_3)
        x = Emu(x + arrow_w + gap)
add_caption(s, Inches(0.8), Inches(5.4), Inches(11.7),
            "From PDF intake to customer email — the same system the whole way through.",
            size=13)
add_footer(s, "03 / 11")


# ── Slide 4 — Live production floor (kanban) ────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "First question of the day, answered without a phone call.",
            kicker="LIVE PRODUCTION FLOOR")
kanban_img = ASSETS / "02-kanban.png"
if kanban_img.exists():
    s.shapes.add_picture(str(kanban_img),
                         Inches(0.6), Inches(1.95),
                         width=Inches(12.1), height=Inches(4.7))
callouts = [
    "Hospital lane — every blocked Job Card surfaces here, regardless of week.",
    "Week badges (W19 / W20) — Job Cards carrying over from earlier weeks are obvious.",
    "Force-advance — supervisor moves a Job Card past a sticking station in 2 clicks.",
]
add_bullets(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5),
            callouts, size=11, line_spacing=1.0)
add_footer(s, "04 / 11")


# ── Slide 5 — Mobile tech experience ───────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "If a tech needs more than two taps, they won't use it.",
            kicker="MOBILE TECH EXPERIENCE")
img_top = Inches(2.0)
img_h   = Inches(4.5)
left_x  = Inches(1.5)
right_x = Inches(4.4)
for img, x in [(ASSETS / "06-tech-list.png", left_x),
               (ASSETS / "07-tech-detail.png", right_x)]:
    if img.exists():
        s.shapes.add_picture(str(img), x, img_top, height=img_h)
text_x = Inches(7.7)
text_w = Inches(5.2)
add_text(s, text_x, Inches(2.1), text_w, Inches(0.6),
         "One tap to clock in.",
         font=DISPLAY_FONT, size=26, color=INK)
add_text(s, text_x, Inches(2.85), text_w, Inches(0.6),
         "Zero PCs on the floor.",
         font=DISPLAY_FONT, size=26, color=INK)
add_bullets(s, text_x, Inches(3.85), text_w, Inches(2.3),
            [
                "Median touchpoint: under 30 seconds per task event.",
                "Blocked Job Cards remain visible — read-only — so the tech sees the queue clearing.",
                "No shared logins. No end-of-day data entry.",
            ], size=14, line_spacing=1.5)
add_footer(s, "05 / 11")


# ── Slide 6 — Insights without a BI team ──────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "The reports the supervisor would have asked for in three months — already there on day one.",
            kicker="INSIGHTS WITHOUT A BI TEAM")
panel_imgs = [
    (ASSETS / "03-variance.png",       "Variance Root Cause",
     "Drilling jig broken accounts for\n38h of overrun this quarter.\nShown automatically."),
    (ASSETS / "04-concentration.png",  "Customer Concentration",
     "Top 3 customers = 67% of hours.\nDFE alone = 42%.\nPareto and 8-quarter trend."),
    (ASSETS / "05-forecast.png",       "Strategic Forecast",
     "Job Card 0007 = HIGH risk.\nPaint overcommitted W21.\n18% avg overrun on TP42N."),
]
panel_w  = Inches(4.0)
panel_gap = Inches(0.3)
panel_top = Inches(2.0)
panel_x = Inches(0.5)
for img, head, body in panel_imgs:
    add_rect(s, panel_x, panel_top, panel_w, Inches(4.7), PAPER)
    if img.exists():
        s.shapes.add_picture(str(img), panel_x + Inches(0.15),
                             panel_top + Inches(0.15),
                             width=panel_w - Inches(0.3))
    cap_top = panel_top + Inches(3.4)
    add_text(s, panel_x + Inches(0.2), cap_top, panel_w - Inches(0.4), Inches(0.4),
             head, font=DISPLAY_FONT, size=15, color=INK)
    add_text(s, panel_x + Inches(0.2), cap_top + Inches(0.45),
             panel_w - Inches(0.4), Inches(1.0),
             body, font=SANS_FONT, size=11, color=INK_3, line_spacing=1.4)
    panel_x = Emu(panel_x + panel_w + panel_gap)
add_footer(s, "06 / 11")


# ── Slide 7 — Excel vs NEE Forge (scenarios) ──────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "Excel is free until you count the reconciliation hours.",
            kicker="EXCEL vs NEE FORGE")
excel_rows = [
    ("\"Where's Job Card 0007 right now?\"",
     "5 phone calls + walk to floor",
     "2-second search; current station + assigned tech visible"),
    ("\"How many hours of overrun on TP42N\nlast quarter?\"",
     "Pivot table + manual reconciliation,\n~30 min — if the data is even there",
     "One click on the Variance Root\nCause report"),
    ("\"Who clocked in on Job 0042 last\nTuesday?\"",
     "Not tracked",
     "Time entries audit log per task"),
    ("\"How many Job Cards from DFE last\nquarter? Year-on-year?\"",
     "Master sheet running total — if anyone\nremembered to update it",
     "Customer Concentration tab\nwith 8-quarter trend"),
    ("\"Tech reports drilling jig broken at 14:30\"",
     "End-of-day note in WhatsApp",
     "Real-time blocker on the kanban;\nsupervisor unblocks with notes"),
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.5),
          ["WORKSHOP SCENARIO", "EXCEL TODAY", "NEE FORGE"],
          excel_rows, font=SANS_FONT, size=11)
add_footer(s, "07 / 11")


# ── Slide 8 — Kinetic vs NEE Forge ─────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "Kinetic is a kit. NEE Forge is the finished cabinet.",
            kicker="EPICOR KINETIC vs NEE FORGE")
kinetic_rows = [
    ("Kanban board",
     "Generic ERP kanban — stations and lanes\nconfigured by external implementer",
     "Body-type tracks (BODY · CHASSIS · SUBFRAME)\nand Hospital lane built in. No config."),
    ("Tech clock-in",
     "Full MES — multi-shift, indirect labour,\nbadge readers",
     "One tap on a phone. No badge reader,\nno shift configuration."),
    ("Variance tracking",
     "Cost-variance categories with multi-level\napproval workflow",
     "13 reason codes attached at task complete.\nRoot-cause aggregate same day."),
    ("QC",
     "Multi-step inspection plans + sample sizes,\nNCR workflow",
     "6-item checklist that triggers a customer-\nfacing email automatically."),
    ("Drafter handoff",
     "Engineering Change Order subsystem\n(custom configuration)",
     "Dedicated drafter persona — Layout / BOM /\nDrawing pack uploads."),
    ("Implementation effort",
     "6 – 18 months. Off-shore implementation.",
     "Days. Shipped as a finished product."),
    ("Modules you'll never use",
     "Multi-site · Multi-currency · MRP · EDI · HR ·\nPayroll · Field Service",
     "None — system is exactly the size\nof the problem."),
]
add_table(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(4.7),
          ["CAPABILITY", "EPICOR KINETIC", "NEE FORGE"],
          kinetic_rows, font=SANS_FONT, size=10)
add_caption(s, Inches(0.5), Inches(6.8), Inches(12.3),
            "Both systems can do most things. Only one is built for this workshop.",
            size=12)
add_footer(s, "08 / 11")


# ── Slide 9 — Implementation timeline ──────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "Three options. Three timelines. Same calendar.",
            kicker="IMPLEMENTATION TIMELINE")
for m_label, x in [("Today", Inches(2.0)), ("M3", Inches(3.7)),
                   ("M6", Inches(5.4)), ("M9", Inches(7.1)),
                   ("M12", Inches(8.8)), ("M15", Inches(10.5)),
                   ("M18", Inches(12.2))]:
    add_text(s, x, Inches(2.05), Inches(0.7), Inches(0.3),
             m_label, font=MONO_FONT, size=9, color=INK_3,
             align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.3), Inches(2.4),
             Emu(int(Inches(0.01))), Inches(0.15), INK_3)

axis_y0 = Inches(2.6)
row_h   = Inches(0.6)
row_gap = Inches(0.3)
left_x  = Inches(2.0)
total_w = Inches(10.4)

def bar(slide, label, frac_start, frac_end, y, color, caption):
    add_text(slide, Inches(0.4), y - Inches(0.05), Inches(1.7), row_h, label,
             font=DISPLAY_FONT, size=15, color=INK)
    add_rect(slide, left_x, y + Inches(0.22),
             total_w, Emu(int(Inches(0.04))), PAPER_3)
    x0 = left_x + Emu(int(total_w * frac_start))
    x1 = left_x + Emu(int(total_w * frac_end))
    add_rect(slide, x0, y, Emu(x1 - x0), row_h, color)
    cap_x = x1 + Inches(0.15)
    add_text(slide, cap_x, y + Inches(0.05), Inches(5), row_h, caption,
             font=SANS_FONT, size=11, color=INK_2)

bar(s, "Excel today",       0.0, 1.0,  axis_y0,                          INK_3,
    "Indefinite ongoing reconciliation cost.")
bar(s, "Epicor Kinetic",    0.0, 1.0,  axis_y0 + row_h + row_gap,        WARN,
    "Industry-standard implementation timeline.")
bar(s, "NEE Forge",         0.0, 0.11, axis_y0 + (row_h + row_gap) * 2,  ACCENT,
    "5-day UAT  ·  30-day evaluation  ·  live in 2 months.")

add_caption(s, Inches(0.5), Inches(6.6), Inches(12.3),
            "Pick a UAT start date this week, and Forge is on a public URL by Friday.",
            size=12)
add_footer(s, "09 / 11")


# ── Slide 10 — Numbers ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s, INK)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
add_text(s, Inches(0.5), Inches(0.5), Inches(8), Inches(0.4),
         "THE NUMBERS THAT CLOSE THE ROOM",
         font=MONO_FONT, size=11, color=ACCENT, bold=True)

boxes = [
    ("5",   "days to a public UAT URL\nwith your supervisor + 2 techs"),
    ("30",  "seconds — median tech touchpoint\n(clock-in / out / complete / blocker)"),
    ("24",  "hours / month of supervisor's\ntime saved versus reconciling\nthree Excel sheets"),
    ("0",   "per-user licence fees.\nPredictable monthly hosting.\nSingle line on the invoice."),
]
cell_w = Inches(6.0)
cell_h = Inches(2.65)
grid_x = Inches(0.65)
grid_y = Inches(1.55)
gap_x  = Inches(0.4)
gap_y  = Inches(0.3)
for i, (num, body) in enumerate(boxes):
    col = i % 2
    row = i // 2
    x = grid_x + Emu(col * (cell_w + gap_x))
    y = grid_y + Emu(row * (cell_h + gap_y))
    add_text(s, x, y, cell_w, Inches(1.6), num,
             font=DISPLAY_FONT, size=160, color=PAPER)
    add_text(s, x, y + Inches(1.7), cell_w, Inches(1.0), body,
             font=SANS_FONT, size=15, color=PAPER_2, line_spacing=1.4)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "Every number above is verifiable. Ask for the evidence appendix.",
         font=MONO_FONT, size=11, color=PAPER_3, align=PP_ALIGN.CENTER)


# ── Slide 11 — Closing ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
header_band(s, "Two outcomes — your call.", kicker="DECISION")

col_w = Inches(5.9)
col_h = Inches(4.3)
left_x  = Inches(0.6)
right_x = Inches(6.85)

add_rect(s, left_x, Inches(2.0), col_w, col_h, PAPER_2)
add_text(s, left_x + Inches(0.4), Inches(2.15), col_w, Inches(0.5),
         "DO NOTHING", font=MONO_FONT, size=11, color=INK_3, bold=True)
nothing_lines = [
    "Spreadsheets remain the master sheet",
    "Variance reasons remain anecdotal",
    "Phone calls remain the supervisor's day-job",
    "Kinetic remains an 18-month maybe",
    "Cost is invisible but real",
]
y_cursor = Inches(2.7)
for line in nothing_lines:
    add_text(s, left_x + Inches(0.4), y_cursor, col_w - Inches(0.6),
             Inches(0.45), "—  " + line,
             font=SANS_FONT, size=14, color=INK_2)
    y_cursor = Emu(y_cursor + Inches(0.55))

add_rect(s, right_x, Inches(2.0), col_w, col_h, INK)
add_rect(s, right_x, Inches(2.0), col_w, Inches(0.06), ACCENT)
add_text(s, right_x + Inches(0.4), Inches(2.15), col_w, Inches(0.5),
         "ADOPT NEE FORGE", font=MONO_FONT, size=11, color=ACCENT, bold=True)
forge_lines = [
    "UAT live in 5 days on a public URL",
    "Every overrun has a reason code by month two",
    "Job Card location answered in 2 seconds",
    "Decision in 30 days — roll-out or walk away",
    "Cost is line-item visible and capped",
]
y_cursor = Inches(2.7)
for line in forge_lines:
    add_text(s, right_x + Inches(0.4), y_cursor, col_w - Inches(0.6),
             Inches(0.45), "→  " + line,
             font=SANS_FONT, size=14, color=PAPER_2)
    y_cursor = Emu(y_cursor + Inches(0.55))

add_rect(s, 0, Inches(6.5), SLIDE_W, Inches(0.7), ACCENT)
add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
         "Pick a UAT start date this week.",
         font=DISPLAY_FONT, size=24, color=PAPER, align=PP_ALIGN.CENTER, bold=True)
add_footer(s, "11 / 11")


OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT_FILE))
print(f"wrote {OUT_FILE}  ({OUT_FILE.stat().st_size // 1024} KB)")
